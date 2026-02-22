"""
文档解析器模块

提供统一的文档解析接口，支持：
- PDF 文档解析
- PPTX 演示文稿解析
- DOCX Word 文档解析
- 网页转 Markdown
"""

from abc import ABC, abstractmethod
from typing import Dict, Any, Optional
from pathlib import Path
import httpx
from modules.YA_Common.utils.logger import get_logger

logger = get_logger("file_parser")


class BaseParser(ABC):
    """文档解析器基类"""
    
    @abstractmethod
    async def parse(self, source: str) -> Dict[str, Any]:
        """
        解析文档
        
        Args:
            source: 文档来源（文件路径或 URL）
            
        Returns:
            包含 content、title、metadata 的字典
        """
        pass
    
    @abstractmethod
    def can_handle(self, source: str) -> bool:
        """判断是否能处理该类型的文档"""
        pass


class PDFParser(BaseParser):
    """PDF 文档解析器，使用 pypdfium2"""
    
    def can_handle(self, source: str) -> bool:
        """检查是否为 PDF 文件"""
        return source.lower().endswith('.pdf') and not source.startswith(('http://', 'https://'))
    
    async def parse(self, source: str) -> Dict[str, Any]:
        """
        解析 PDF 文件
        
        Args:
            source: PDF 文件路径
            
        Returns:
            {
                "content": str,  # 提取的文本内容
                "title": str,    # 文档标题（从文件名提取）
                "metadata": {    # 元数据
                    "pages": int,
                    "file_path": str,
                    "doc_type": "pdf"
                }
            }
        """
        try:
            import pypdfium2 as pdfium
        except ImportError:
            raise RuntimeError(
                "需要安装 pypdfium2 库: pip install pypdfium2 (或 uv add pypdfium2)"
            )
        
        # 清理文件路径
        source = source.strip()
        path = Path(source)
        if not path.exists():
            raise FileNotFoundError(f"PDF 文件不存在: {source}")
        
        logger.info(f"开始解析 PDF: {source}")
        
        try:
            # 使用 pypdfium2 读取 PDF
            pdf = pdfium.PdfDocument(str(path))
            total_pages = len(pdf)
            
            # 提取所有页面文本
            all_text = []
            for page_num in range(total_pages):
                page = pdf[page_num]
                textpage = page.get_textpage()
                text = textpage.get_text_range()
                
                if text.strip():
                    all_text.append(f"## 第 {page_num + 1} 页\n\n{text.strip()}")
            
            content = "\n\n".join(all_text)
            
            # 从文件名提取标题
            title = path.stem.replace('_', ' ').replace('-', ' ')
            
            logger.info(f"PDF 解析完成: {total_pages} 页，{len(content)} 字符")
            
            return {
                "content": content,
                "title": title,
                "metadata": {
                    "pages": total_pages,
                    "file_path": str(path.absolute()),
                    "doc_type": "pdf"
                }
            }
        except Exception as e:
            logger.error(f"PDF 解析失败: {e}")
            raise RuntimeError(f"PDF 解析失败: {str(e)}")


class PPTXParser(BaseParser):
    """PPTX 演示文稿解析器，使用 python-pptx"""

    def can_handle(self, source: str) -> bool:
        """检查是否为 PPTX 文件"""
        return source.lower().endswith('.pptx') and not source.startswith(('http://', 'https://'))

    async def parse(self, source: str) -> Dict[str, Any]:
        """
        解析 PPTX 文件

        Args:
            source: PPTX 文件路径

        Returns:
            {
                "content": str,  # 提取的文本内容（Markdown 格式）
                "title": str,    # 文档标题（从文件名提取）
                "metadata": {    # 元数据
                    "slides": int,
                    "file_path": str,
                    "doc_type": "pptx"
                }
            }
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError(
                "需要安装 python-pptx 库: pip install python-pptx (或 uv add python-pptx)"
            )

        # 清理文件路径
        source = source.strip()
        path = Path(source)
        if not path.exists():
            raise FileNotFoundError(f"PPTX 文件不存在: {source}")

        logger.info(f"开始解析 PPTX: {source}")

        try:
            prs = Presentation(str(path))
            total_slides = len(prs.slides)

            all_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []

                for shape in slide.shapes:
                    # 提取文本框内容
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                # 根据段落层级生成 Markdown
                                level = paragraph.level
                                if level == 0 and slide_texts == []:
                                    slide_texts.append(f"### {text}")
                                else:
                                    indent = "  " * level
                                    slide_texts.append(f"{indent}- {text}")

                    # 提取表格内容
                    if shape.has_table:
                        table = shape.table
                        table_md = self._table_to_markdown(table)
                        if table_md:
                            slide_texts.append(table_md)

                if slide_texts:
                    slide_content = "\n".join(slide_texts)
                    all_text.append(f"## 第 {slide_num} 页\n\n{slide_content}")

            content = "\n\n".join(all_text)

            # 从文件名提取标题
            title = path.stem.replace('_', ' ').replace('-', ' ')

            logger.info(f"PPTX 解析完成: {total_slides} 页，{len(content)} 字符")

            return {
                "content": content,
                "title": title,
                "metadata": {
                    "slides": total_slides,
                    "file_path": str(path.absolute()),
                    "doc_type": "pptx"
                }
            }
        except Exception as e:
            logger.error(f"PPTX 解析失败: {e}")
            raise RuntimeError(f"PPTX 解析失败: {str(e)}")

    def _table_to_markdown(self, table) -> str:
        """将 PPTX 表格转换为 Markdown 表格"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")

        if len(rows) >= 1:
            # 添加表头分隔行
            col_count = len(table.rows[0].cells)
            separator = "| " + " | ".join(["---"] * col_count) + " |"
            rows.insert(1, separator)

        return "\n".join(rows)


class DOCXParser(BaseParser):
    """DOCX Word 文档解析器，使用 python-docx"""

    def can_handle(self, source: str) -> bool:
        """检查是否为 DOCX 文件"""
        return source.lower().endswith('.docx') and not source.startswith(('http://', 'https://'))

    async def parse(self, source: str) -> Dict[str, Any]:
        """
        解析 DOCX 文件

        Args:
            source: DOCX 文件路径

        Returns:
            {
                "content": str,  # 提取的文本内容（Markdown 格式）
                "title": str,    # 文档标题（从文件名或标题样式提取）
                "metadata": {    # 元数据
                    "paragraphs": int,
                    "tables": int,
                    "file_path": str,
                    "doc_type": "docx"
                }
            }
        """
        try:
            from docx import Document
        except ImportError:
            raise RuntimeError(
                "需要安装 python-docx 库: pip install python-docx (或 uv add python-docx)"
            )

        # 清理文件路径
        source = source.strip()
        path = Path(source)
        if not path.exists():
            raise FileNotFoundError(f"DOCX 文件不存在: {source}")

        logger.info(f"开始解析 DOCX: {source}")

        try:
            doc = Document(str(path))

            all_text = []
            doc_title = ""
            paragraph_count = 0
            table_count = len(doc.tables)

            # 遍历文档主体中的所有元素（保持段落和表格的顺序）
            for element in doc.element.body:
                tag = element.tag.split('}')[-1]  # 去除命名空间

                if tag == 'p':  # 段落
                    para = None
                    for p in doc.paragraphs:
                        if p._element is element:
                            para = p
                            break
                    if para is None:
                        continue

                    text = para.text.strip()
                    if not text:
                        continue

                    paragraph_count += 1
                    style_name = para.style.name if para.style else ""

                    # 根据样式转换为 Markdown
                    if style_name.startswith('Heading 1') or style_name == 'Title':
                        all_text.append(f"# {text}")
                        if not doc_title:
                            doc_title = text
                    elif style_name.startswith('Heading 2'):
                        all_text.append(f"## {text}")
                    elif style_name.startswith('Heading 3'):
                        all_text.append(f"### {text}")
                    elif style_name.startswith('Heading 4'):
                        all_text.append(f"#### {text}")
                    elif style_name.startswith('Heading'):
                        all_text.append(f"##### {text}")
                    elif style_name.startswith('List'):
                        all_text.append(f"- {text}")
                    else:
                        all_text.append(text)

                elif tag == 'tbl':  # 表格
                    tbl = None
                    for t in doc.tables:
                        if t._element is element:
                            tbl = t
                            break
                    if tbl is None:
                        continue

                    table_md = self._table_to_markdown(tbl)
                    if table_md:
                        all_text.append(table_md)

            content = "\n\n".join(all_text)

            # 使用文档中的标题或文件名
            if not doc_title:
                doc_title = path.stem.replace('_', ' ').replace('-', ' ')

            logger.info(
                f"DOCX 解析完成: {paragraph_count} 段落, {table_count} 表格, "
                f"{len(content)} 字符"
            )

            return {
                "content": content,
                "title": doc_title,
                "metadata": {
                    "paragraphs": paragraph_count,
                    "tables": table_count,
                    "file_path": str(path.absolute()),
                    "doc_type": "docx"
                }
            }
        except Exception as e:
            logger.error(f"DOCX 解析失败: {e}")
            raise RuntimeError(f"DOCX 解析失败: {str(e)}")

    def _table_to_markdown(self, table) -> str:
        """将 DOCX 表格转换为 Markdown 表格"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")

        if len(rows) >= 1:
            col_count = len(table.rows[0].cells)
            separator = "| " + " | ".join(["---"] * col_count) + " |"
            rows.insert(1, separator)

        return "\n".join(rows)


class WebPageParser(BaseParser):
    """网页转 Markdown 解析器"""
    
    def can_handle(self, source: str) -> bool:
        """检查是否为 URL"""
        return source.startswith(('http://', 'https://'))
    
    async def parse(self, source: str) -> Dict[str, Any]:
        """
        抓取网页并转换为 Markdown
        
        Args:
            source: 网页 URL
            
        Returns:
            {
                "content": str,  # Markdown 格式内容
                "title": str,    # 网页标题
                "metadata": {    # 元数据
                    "url": str,
                    "doc_type": "webpage"
                }
            }
        """
        try:
            import html2text
        except ImportError:
            raise RuntimeError(
                "需要安装 html2text 库: pip install html2text (或 uv add html2text)"
            )
        
        # 清理 URL（去除空格、换行符等）
        source = source.strip()
        
        logger.info(f"开始抓取网页: {source}")
        
        try:
            # 抓取网页内容
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
            }
            
            # 创建客户端，测试环境可能需要禁用 SSL 验证
            async with httpx.AsyncClient(
                timeout=15, 
                follow_redirects=True,
                verify=False  # 生产环境中应设为 True
            ) as client:
                response = await client.get(source, headers=headers)
                response.raise_for_status()
                html_content = response.text
            
            # 提取标题
            title = self._extract_title(html_content, source)
            
            # 转换为 Markdown
            h2t = html2text.HTML2Text()
            h2t.ignore_links = False
            h2t.ignore_images = False
            h2t.ignore_emphasis = False
            h2t.body_width = 0  # 不自动换行
            h2t.mark_code = True
            
            markdown_content = h2t.handle(html_content)
            
            logger.info(f"网页转换完成: {len(markdown_content)} 字符")
            
            return {
                "content": markdown_content,
                "title": title,
                "metadata": {
                    "url": source,
                    "doc_type": "webpage",
                    "html_content": html_content  # 保存原始 HTML 以便后续保存
                }
            }
        except httpx.HTTPError as e:
            logger.error(f"网页抓取失败: {e}")
            raise RuntimeError(f"网页抓取失败: {str(e)}")
        except Exception as e:
            logger.error(f"网页解析失败: {e}")
            raise RuntimeError(f"网页解析失败: {str(e)}")
    
    def _extract_title(self, html: str, fallback_url: str) -> str:
        """从 HTML 中提取标题"""
        try:
            from bs4 import BeautifulSoup
            
            soup = BeautifulSoup(html, "html.parser")
            
            # 优先级：<title> > <h1> > URL
            if soup.title and soup.title.string:
                return soup.title.string.strip()
            
            h1 = soup.find("h1")
            if h1 and h1.get_text():
                return h1.get_text().strip()
            
            # 从 URL 提取
            from urllib.parse import urlparse
            parsed = urlparse(fallback_url)
            return parsed.netloc
            
        except Exception:
            return fallback_url


class DocumentParserFactory:
    """文档解析器工厂类"""
    
    def __init__(self):
        self.parsers = [
            PDFParser(),
            PPTXParser(),
            DOCXParser(),
            WebPageParser(),
        ]
    
    async def parse(self, source: str) -> Dict[str, Any]:
        """
        自动选择合适的解析器解析文档
        
        Args:
            source: 文档来源（文件路径或 URL）
            
        Returns:
            解析结果字典
            
        Raises:
            ValueError: 如果没有合适的解析器
        """
        for parser in self.parsers:
            if parser.can_handle(source):
                logger.info(f"使用解析器: {parser.__class__.__name__}")
                return await parser.parse(source)
        
        raise ValueError(
            f"不支持的文档类型: {source}\n"
            f"当前支持: PDF (.pdf)、PPTX (.pptx)、DOCX (.docx) 和网页 (http/https URL)"
        )
